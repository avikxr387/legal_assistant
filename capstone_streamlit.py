import streamlit as st
import uuid
from graph.graph_builder import build_graph
from rag.vectordb import add_uploaded_doc


# =========================
# TEXT EXTRACTION
# =========================
def extract_text(file):
    name = file.name.lower()

    if name.endswith(".txt"):
        return file.read().decode("utf-8")

    elif name.endswith(".pdf"):
        from pypdf import PdfReader
        reader = PdfReader(file)
        return " ".join([page.extract_text() or "" for page in reader.pages])

    elif name.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return " ".join(text)

    return ""


# =========================
# CACHE HEAVY OBJECTS
# =========================
@st.cache_resource
def load_app():
    return build_graph()

app, collection, model = load_app()


# =========================
# SESSION STATE INIT
# =========================
if "messages" not in st.session_state:
    st.session_state.messages = []

if "thread_id" not in st.session_state:
    st.session_state.thread_id = str(uuid.uuid4())

if "full_doc_text" not in st.session_state:
    st.session_state.full_doc_text = ""


# =========================
# SIDEBAR
# =========================
st.sidebar.title("⚖️ Legal Assistant")

st.sidebar.markdown("""
### Domain
Legal Document Assistant

### Capabilities
- Contract law Q&A  
- Clause explanations  
- Legal concept breakdown  
- Memory-based conversation  

### Topics Covered
- Breach of Contract  
- Damages  
- Termination Clause  
- Consideration  
- Confidentiality  

---

### About
This assistant uses Retrieval-Augmented Generation (RAG), memory, and self-evaluation to provide grounded legal answers.
""")

if st.sidebar.button("🆕 New Conversation"):
    st.session_state.messages = []
    st.session_state.thread_id = str(uuid.uuid4())
    st.session_state.full_doc_text = ""
    collection.delete(where={"topic": "User Upload"})
    st.rerun()


# =========================
# MAIN UI
# =========================
st.title("📜 Legal Document Assistant")


# =========================
# FILE UPLOAD (CORRECT POSITION)
# =========================
uploaded_file = st.file_uploader(
    "Upload document (txt, pdf, pptx)",
    type=["txt", "pdf", "pptx"]
)

if uploaded_file:
    text = extract_text(uploaded_file)

    # 🔥 validation
    if not text or len(text.strip()) < 50:
        st.error("Document too small or unreadable.")
    else:
        st.session_state.full_doc_text = text

        add_uploaded_doc(collection, model, text)

        # 🔥 UX improvements
        st.success(f"Uploaded: {uploaded_file.name}")
        st.info(f"Document length: {len(text)} characters")
        st.info(f"Chunks created: {len(text) // 400}")
        
if uploaded_file and uploaded_file.size > 5 * 1024 * 1024:
    st.error("File too large. Max 5MB allowed.")
    st.stop()


# =========================
# DISPLAY CHAT HISTORY
# =========================
for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])


# =========================
# USER INPUT
# =========================
user_input = st.chat_input("Ask a legal question...")

if user_input:
    # Show user message
    st.session_state.messages.append({"role": "user", "content": user_input})

    with st.chat_message("user"):
        st.markdown(user_input)

    # Prepare state
    state = {
        "question": user_input,
        "route": "",
        "retrieved": "",
        "sources": [],
        "tool_result": None,
        "answer": "",
        "faithfulness": 0.0,
        "eval_retries": 0,
        "full_doc_text": st.session_state.full_doc_text
    }

    # Run graph
    result = app.invoke(
        state,
        config={"configurable": {"thread_id": st.session_state.thread_id}}
    )

    answer = result["answer"]

    # Add sources
    if result.get("sources"):
        answer += "\n\n**Sources:** " + ", ".join(result["sources"])

    # Add faithfulness (optional but impressive)
    if "faithfulness" in result:
        answer += f"\n\n*(Faithfulness: {round(result['faithfulness'], 2)})*"

    # Show assistant response
    with st.chat_message("assistant"):
        st.markdown(answer)

    # Save response
    st.session_state.messages.append({"role": "assistant", "content": answer})