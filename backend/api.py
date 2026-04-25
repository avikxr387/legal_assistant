import os
import warnings
import logging

os.environ["TRANSFORMERS_VERBOSITY"] = "error"
os.environ["TOKENIZERS_PARALLELISM"] = "false"
os.environ["HF_HUB_DISABLE_SYMLINKS_WARNING"] = "1"

warnings.filterwarnings("ignore")

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s"
)

logging.getLogger("sentence_transformers").setLevel(logging.ERROR)


from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
import uuid
import io

from graph.graph_builder import build_graph
from rag.vectordb import add_uploaded_doc

from pypdf import PdfReader
from pptx import Presentation


# =========================
# INIT APP
# =========================
app = FastAPI(title="Legal Assistant API")


# =========================
# CORS (allow frontend)
# =========================
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # change in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# =========================
# LOAD GRAPH (ONCE)
# =========================
app_graph, collection, model = build_graph()


# =========================
# SESSION STORE (IN-MEMORY)
# =========================
sessions = {}  # thread_id -> document text


# =========================
# HELPER: EXTRACT TEXT
# =========================
def extract_text(filename: str, content: bytes) -> str:
    filename = filename.lower()

    if filename.endswith(".txt"):
        return content.decode("utf-8", errors="ignore")

    elif filename.endswith(".pdf"):
        reader = PdfReader(io.BytesIO(content))
        return "\n".join([p.extract_text() or "" for p in reader.pages])

    elif filename.endswith(".pptx"):
        prs = Presentation(io.BytesIO(content))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    else:
        raise HTTPException(status_code=400, detail="Unsupported file type")


# =========================
# ROOT
# =========================
@app.get("/")
def root():
    return {"message": "Legal Assistant API running"}


# =========================
# CREATE SESSION
# =========================
@app.post("/session")
def create_session():
    thread_id = str(uuid.uuid4())
    sessions[thread_id] = ""
    return {"thread_id": thread_id}


# =========================
# FILE UPLOAD
# =========================
@app.post("/upload")
async def upload_file(thread_id: str, file: UploadFile = File(...)):
    if thread_id not in sessions:
        raise HTTPException(status_code=400, detail="Invalid thread_id")

    content = await file.read()

    text = extract_text(file.filename, content)

    if not text.strip():
        raise HTTPException(status_code=400, detail="Empty document")

    # store full text in session
    sessions[thread_id] = text

    # add to vector DB
    add_uploaded_doc(collection, model, text)

    return {
        "message": "Document uploaded successfully",
        "length": len(text)
    }


# =========================
# CHAT
# =========================
@app.post("/chat")
async def chat(data: dict):
    question = data.get("question")
    thread_id = data.get("thread_id")

    if not question:
        raise HTTPException(status_code=400, detail="Question required")

    if thread_id not in sessions:
        raise HTTPException(status_code=400, detail="Invalid thread_id")

    # build state
    state = {
        "question": question,
        "route": "",
        "retrieved": "",
        "sources": [],
        "tool_result": None,
        "answer": "",
        "faithfulness": 0.0,
        "eval_retries": 0,
        "full_doc_text": sessions[thread_id]
    }

    # run graph
    result = app_graph.invoke(
        state,
        config={"configurable": {"thread_id": thread_id}}
    )

    return {
        "answer": result.get("answer"),
        "sources": result.get("sources", []),
        "faithfulness": result.get("faithfulness", 0.0)
    }